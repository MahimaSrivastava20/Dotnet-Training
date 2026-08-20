from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color Palette ──────────────────────────────────────────────
DEEP_NAVY   = RGBColor(0x0A, 0x0E, 0x27)
VIVID_BLUE  = RGBColor(0x00, 0x5B, 0xBF)
ELECTRIC    = RGBColor(0x00, 0xC2, 0xFF)
VIOLET      = RGBColor(0x7C, 0x3A, 0xED)
TEAL        = RGBColor(0x00, 0xD4, 0xAA)
ROSE        = RGBColor(0xFF, 0x4D, 0x6D)
AMBER       = RGBColor(0xFF, 0xBF, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE2, 0xE8, 0xF0)
MID_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD   = RGBColor(0x10, 0x17, 0x35)
CARD_BORDER = RGBColor(0x1E, 0x2A, 0x4A)

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w if line_w else Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bg(slide, color=DEEP_NAVY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def accent_bar(slide, color=VIVID_BLUE, h=0.06):
    add_rect(slide, 0, 0, 13.33, h, fill=color)
    add_rect(slide, 0, 7.5-h, 13.33, h, fill=color)

def card(slide, l, t, w, h, fill=DARK_CARD, border=CARD_BORDER):
    add_rect(slide, l, t, w, h, fill=fill, line=border, line_w=Pt(1.5))

def section_header(slide, title, subtitle="", bar_color=VIVID_BLUE):
    bg(slide)
    add_rect(slide, 0, 0, 13.33, 1.1, fill=bar_color)
    add_text(slide, title, 0.4, 0.15, 12, 0.8, size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.9, 12, 0.5, size=14, color=LIGHT_GRAY, align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 1.1, 13.33, 0.04, fill=ELECTRIC)

blank = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
# gradient strip left
add_rect(s, 0, 0, 4.5, 7.5, fill=RGBColor(0x05, 0x0A, 0x1E))
# accent diagonal strip
add_rect(s, 0, 0, 0.3, 7.5, fill=VIVID_BLUE)
add_rect(s, 0.3, 0, 0.08, 7.5, fill=ELECTRIC)

add_text(s, "SafeHaven", 0.7, 1.0, 4.0, 0.7, size=14, bold=True, color=ELECTRIC)
add_text(s, "Insurance Policy &\nClaims Support System", 0.7, 1.6, 5.8, 2.2, size=34, bold=True, color=WHITE)
add_rect(s, 0.7, 3.85, 3.5, 0.05, fill=VIVID_BLUE)
add_text(s, "Microservices · Angular 19 · .NET 8 · RabbitMQ", 0.7, 4.0, 6.0, 0.5, size=13, color=MID_GRAY)
add_text(s, "Internship Project Presentation", 0.7, 6.6, 5.0, 0.5, size=12, color=MID_GRAY, italic=True)

# right side visual cards
colors = [VIVID_BLUE, VIOLET, TEAL, ROSE, AMBER, ELECTRIC]
labels = ["Identity\nService", "Policy\nService", "Ticket\nService", "Payment\nService", "Notification\nService", "Admin\nService"]
positions = [(5.6,0.5),(7.5,0.5),(9.4,0.5),(5.6,2.8),(7.5,2.8),(9.4,2.8)]
for i,(lx,ty) in enumerate(positions):
    card(s, lx, ty, 1.8, 2.1, fill=RGBColor(int(colors[i][0]*0.3)//1,int(colors[i][1]*0.3)//1,int(colors[i][2]*0.3)//1), border=colors[i])
    add_rect(s, lx, ty, 1.8, 0.35, fill=colors[i])
    add_text(s, labels[i], lx+0.1, ty+0.4, 1.6, 1.5, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "API Gateway (Ocelot)  ·  SQL Server  ·  Razorpay  ·  JWT Auth",
         5.4, 5.2, 7.5, 0.5, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "Agenda", "What we'll cover today", bar_color=VIOLET)

items = [
    ("01", "Tech Stack",          "Languages, frameworks, tools & databases",        VIVID_BLUE),
    ("02", "System Architecture", "Microservices layout and API Gateway",             VIOLET),
    ("03", "Project Flow",        "End-to-end user journey from login to payment",    TEAL),
    ("04", "Key Services",        "Deep-dive into each microservice",                 ROSE),
    ("05", "Frontend Pages",      "Angular 19 UI, routing, and components",           AMBER),
    ("06", "Security & Events",   "JWT authentication and RabbitMQ messaging",        ELECTRIC),
]
cols = [(0.3, 1.3), (6.8, 1.3)]
for i, (num, title, desc, col) in enumerate(items):
    cx, cy = cols[i % 2]
    row = i // 2
    ty = cy + row * 2.0
    card(s, cx, ty, 6.2, 1.7, fill=DARK_CARD, border=col)
    add_rect(s, cx, ty, 0.55, 1.7, fill=col)
    add_text(s, num, cx+0.05, ty+0.5, 0.5, 0.7, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, cx+0.65, ty+0.15, 5.4, 0.5, size=16, bold=True, color=WHITE)
    add_text(s, desc,  cx+0.65, ty+0.65, 5.3, 0.8, size=11, color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — TECH STACK
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "Tech Stack", "Every technology powering SafeHaven", bar_color=VIVID_BLUE)

tech = [
    ("Frontend",  VIVID_BLUE, ["Angular 19 (Standalone)", "TypeScript 5", "TailwindCSS", "RxJS", "Angular Router"]),
    ("Backend",   VIOLET,     [".NET 8 / ASP.NET Core", "Entity Framework Core", "RESTful APIs", "Ocelot API Gateway", "CORS & Middleware"]),
    ("Database",  TEAL,       ["SQL Server 2022", "EF Core Migrations", "Code-First Models", "Relational Schema", "Multi-DB (per service)"]),
    ("Messaging", ROSE,       ["RabbitMQ", "Event-Driven Arch.", "Publisher / Consumer", "Queue per domain", "Async Notifications"]),
    ("Security",  AMBER,      ["JWT Bearer Tokens", "Role-Based Auth", "bcrypt Password Hash", "HTTPS / SSL", "Admin & User Roles"]),
    ("Payment",   ELECTRIC,   ["Razorpay SDK", "Order Creation API", "Payment Verification", "Signature Validation", "INR Currency"]),
]
for i, (name, col, items) in enumerate(tech):
    cx = 0.25 + (i % 3) * 4.33
    ty = 1.25 + (i // 3) * 3.0
    card(s, cx, ty, 4.1, 2.7, fill=DARK_CARD, border=col)
    add_rect(s, cx, ty, 4.1, 0.45, fill=col)
    add_text(s, name, cx+0.15, ty+0.05, 3.8, 0.4, size=15, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(s, f"▸  {item}", cx+0.15, ty+0.55+j*0.42, 3.8, 0.4, size=10.5, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "System Architecture", "Microservices with Ocelot API Gateway", bar_color=TEAL)

# Angular box
card(s, 0.3, 1.25, 2.2, 1.1, fill=RGBColor(0x00,0x30,0x6B), border=VIVID_BLUE)
add_text(s, "Angular 19\nFrontend", 0.3, 1.25, 2.2, 1.1, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow → Gateway
add_rect(s, 2.55, 1.65, 1.0, 0.08, fill=ELECTRIC)
add_text(s, "HTTP", 2.6, 1.45, 1.0, 0.3, size=9, color=ELECTRIC, align=PP_ALIGN.CENTER)

# Gateway box
card(s, 3.6, 1.25, 2.4, 1.1, fill=RGBColor(0x1E,0x0A,0x35), border=VIOLET)
add_rect(s, 3.6, 1.25, 2.4, 0.35, fill=VIOLET)
add_text(s, "Ocelot API Gateway", 3.65, 1.28, 2.3, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "JWT Validation\nRoute Forwarding", 3.65, 1.65, 2.3, 0.65, size=9.5, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Arrow → services
add_rect(s, 6.05, 1.68, 0.65, 0.08, fill=ELECTRIC)

# Services column
svc_data = [
    ("Identity Service",      VIVID_BLUE, 6.75, 1.2),
    ("Policy Service",        VIOLET,     6.75, 2.25),
    ("Ticket Service",        ROSE,       6.75, 3.3),
    ("Payment Service",       AMBER,      6.75, 4.35),
    ("Notification Service",  TEAL,       6.75, 5.4),
    ("Admin Service",         ELECTRIC,   6.75, 6.0),
]
for name, col, lx, ty in svc_data:
    card(s, lx, ty, 2.3, 0.75, fill=DARK_CARD, border=col)
    add_rect(s, lx, ty, 0.18, 0.75, fill=col)
    add_text(s, name, lx+0.25, ty+0.18, 2.0, 0.4, size=10, bold=True, color=WHITE)
    # connecting line from gateway
    add_rect(s, 6.05, ty+0.34, 0.72, 0.05, fill=col)

# SQL Server
card(s, 9.25, 1.8, 2.0, 1.0, fill=DARK_CARD, border=TEAL)
add_text(s, "SQL Server\n(per service DB)", 9.25, 1.8, 2.0, 1.0, size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_rect(s, 9.1, 2.28, 0.17, 0.05, fill=TEAL)

# RabbitMQ
card(s, 9.25, 3.1, 2.0, 1.0, fill=DARK_CARD, border=ROSE)
add_text(s, "RabbitMQ\nMessage Broker", 9.25, 3.1, 2.0, 1.0, size=11, bold=True, color=ROSE, align=PP_ALIGN.CENTER)
add_rect(s, 9.1, 3.58, 0.17, 0.05, fill=ROSE)

# Razorpay
card(s, 9.25, 4.4, 2.0, 0.9, fill=DARK_CARD, border=AMBER)
add_text(s, "Razorpay\nPayment Gateway", 9.25, 4.4, 2.0, 0.9, size=11, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — PROJECT FLOW
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "Project Flow", "End-to-end user journey inside SafeHaven", bar_color=ROSE)

steps = [
    ("1", "Register /\nLogin",        VIVID_BLUE),
    ("2", "Browse\nPolicies",         VIOLET),
    ("3", "Select &\nSubscribe",      TEAL),
    ("4", "Razorpay\nPayment",        AMBER),
    ("5", "Policy\nActivated",        ELECTRIC),
    ("6", "Submit\nClaim",            ROSE),
    ("7", "Admin\nReview",            VIVID_BLUE),
    ("8", "Notification\nSent",       TEAL),
]
box_w, box_h = 1.35, 1.5
gap = 0.25
start_x = 0.3
ty = 1.9
for i, (num, label, col) in enumerate(steps):
    lx = start_x + i*(box_w+gap)
    card(s, lx, ty, box_w, box_h, fill=DARK_CARD, border=col)
    add_rect(s, lx, ty, box_w, 0.38, fill=col)
    add_text(s, num, lx, ty+0.04, box_w, 0.3, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, lx+0.05, ty+0.5, box_w-0.1, 0.9, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        ax = lx + box_w + 0.02
        add_rect(s, ax, ty+box_h/2-0.03, gap+0.02, 0.06, fill=col)

# description row
descs = [
    ("JWT token issued,\nroles assigned", VIVID_BLUE),
    ("View Life, Health,\nCar, Home plans", VIOLET),
    ("Choose coverage\namount & term", TEAL),
    ("Secure checkout,\norder verified", AMBER),
    ("Policy stored,\nstatus = Active", ELECTRIC),
    ("Upload evidence,\nclaim amount", ROSE),
    ("Approve / Reject\nin Admin panel", VIVID_BLUE),
    ("Email/in-app via\nRabbitMQ event", TEAL),
]
for i, (desc, col) in enumerate(descs):
    lx = start_x + i*(box_w+gap)
    add_text(s, desc, lx, ty+box_h+0.15, box_w, 0.9, size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — KEY MICROSERVICES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "Key Microservices", "Independent services each with their own DB", bar_color=VIOLET)

services = [
    ("Identity Service",     VIVID_BLUE, "Register, Login, JWT issuance\nRole management (Admin / Customer)\nEmail verification flow"),
    ("Policy Service",       VIOLET,     "CRUD for insurance policies\nLife, Health, Car, Home types\nCoverage & premium tracking"),
    ("Ticket / Claim",       ROSE,       "Claim submission with evidence\nAdmin approval workflow\nStatus: Pending → Approved/Rejected"),
    ("Payment Service",      AMBER,      "Razorpay order creation\nSignature verification\nPayment history per customer"),
    ("Notification Service", TEAL,       "Consumes RabbitMQ events\nSends alerts on claim & payment\nIn-app + email style notifications"),
    ("Admin Service",        ELECTRIC,   "Dashboard: users, policies, stats\nClaim & ticket management\nSystem-wide reporting"),
]
for i, (name, col, desc) in enumerate(services):
    cx = 0.25 + (i % 3) * 4.33
    ty = 1.3 + (i // 3) * 2.9
    card(s, cx, ty, 4.1, 2.6, fill=DARK_CARD, border=col)
    add_rect(s, cx, ty, 4.1, 0.45, fill=col)
    add_text(s, name, cx+0.15, ty+0.06, 3.8, 0.35, size=14, bold=True, color=WHITE)
    add_text(s, desc, cx+0.15, ty+0.55, 3.8, 1.9, size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — FRONTEND PAGES
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "Frontend — Angular 19", "Lazy-loaded standalone components with auth guards", bar_color=AMBER)

pages = [
    ("/",              "Welcome Page",         "Landing with CTA",                 VIVID_BLUE),
    ("/login",         "Login",                "JWT auth + role redirect",          VIOLET),
    ("/register",      "Register",             "User sign-up + verify email",       TEAL),
    ("/dashboard",     "Customer Dashboard",   "My policies, claims summary",       ROSE),
    ("/policies",      "Policy Management",    "Browse & subscribe to policies",    AMBER),
    ("/payment",       "Payment Page",         "Razorpay checkout integration",     ELECTRIC),
    ("/claims",        "Claims Review",        "View all submitted claims",         VIVID_BLUE),
    ("/submit-claim",  "Submit Claim",         "Form with evidence upload",         VIOLET),
    ("/support",       "Support / Helpp",      "Raise support tickets",             TEAL),
    ("/admin",         "Admin Dashboard",      "Manage users, claims, reports",     ROSE),
]
for i, (route, name, desc, col) in enumerate(pages):
    cx = 0.25 + (i % 2) * 6.55
    ty = 1.25 + (i // 2) * 0.62
    card(s, cx, ty, 6.3, 0.55, fill=DARK_CARD, border=col)
    add_rect(s, cx, ty, 0.12, 0.55, fill=col)
    add_text(s, route, cx+0.2, ty+0.1, 2.0, 0.35, size=10, bold=True, color=col)
    add_text(s, name,  cx+2.25, ty+0.1, 2.0, 0.35, size=10, bold=True, color=WHITE)
    add_text(s, desc,  cx+4.3, ty+0.1, 1.9, 0.35, size=9,  color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — SECURITY & RABBITMQ
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
section_header(s, "Security & Event-Driven Messaging", "JWT Auth + RabbitMQ async communication", bar_color=ELECTRIC)

# JWT side
card(s, 0.3, 1.3, 6.1, 5.8, fill=DARK_CARD, border=VIVID_BLUE)
add_rect(s, 0.3, 1.3, 6.1, 0.5, fill=VIVID_BLUE)
add_text(s, "🔐  JWT Authentication", 0.5, 1.33, 5.8, 0.42, size=16, bold=True, color=WHITE)
jwt_points = [
    "User logs in → Identity Service validates credentials",
    "JWT token generated with userId, email, role (Admin/Customer)",
    "Token sent in Authorization: Bearer <token> header",
    "Ocelot Gateway validates token before forwarding requests",
    "Angular stores token in localStorage, attaches via interceptor",
    "Token expiry handled — user redirected to login on 401",
    "Roles enforced: adminGuard blocks non-admin from /admin route",
]
for j, pt in enumerate(jwt_points):
    add_text(s, f"▸  {pt}", 0.5, 1.95+j*0.68, 5.7, 0.6, size=11, color=LIGHT_GRAY)

# RabbitMQ side
card(s, 6.75, 1.3, 6.25, 5.8, fill=DARK_CARD, border=ROSE)
add_rect(s, 6.75, 1.3, 6.25, 0.5, fill=ROSE)
add_text(s, "🐰  RabbitMQ Messaging", 6.95, 1.33, 5.9, 0.42, size=16, bold=True, color=WHITE)
rmq_points = [
    "Services publish events without direct coupling",
    "Payment success → publishes PolicyActivated event",
    "Claim submitted → publishes ClaimReceived event",
    "Notification Service subscribes to all domain queues",
    "Sends notifications on policy activation & claim update",
    "Admin Service consumes events for real-time reporting",
    "Decoupled design → services can scale independently",
]
for j, pt in enumerate(rmq_points):
    add_text(s, f"▸  {pt}", 6.95, 1.95+j*0.68, 5.9, 0.6, size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — CONCLUSION
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
add_rect(s, 0, 0, 13.33, 0.08, fill=TEAL)
add_rect(s, 0, 7.42, 13.33, 0.08, fill=TEAL)

add_text(s, "Key Learnings & Takeaways", 0.5, 0.3, 12, 0.7, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 4.5, 0.95, 4.33, 0.05, fill=TEAL)

learnings = [
    ("Microservices Design",    VIVID_BLUE, "Built 6 independently deployable services each with their own database and responsibility"),
    ("Event-Driven Architecture", ROSE,    "Used RabbitMQ to decouple services — notifications fire without tight coupling"),
    ("Secure Auth",             VIOLET,    "Implemented JWT with role-based access and Ocelot Gateway validation at the edge"),
    ("Full-Stack Integration",  TEAL,      "Connected Angular 19 frontend to .NET 8 APIs seamlessly using HTTP interceptors"),
    ("Payment Integration",     AMBER,     "Integrated Razorpay with order creation, signature verification, and status tracking"),
    ("Real-World Problem",      ELECTRIC,  "Solved a tangible insurance domain problem — policies, claims, payments, and support"),
]
for i, (title, col, desc) in enumerate(learnings):
    cx = 0.3 + (i%2)*6.55
    ty = 1.15 + (i//2)*1.95
    card(s, cx, ty, 6.3, 1.75, fill=DARK_CARD, border=col)
    add_rect(s, cx, ty, 6.3, 0.38, fill=col)
    add_text(s, title, cx+0.15, ty+0.05, 5.9, 0.3, size=13, bold=True, color=WHITE)
    add_text(s, desc, cx+0.15, ty+0.48, 5.9, 1.1, size=10.5, color=LIGHT_GRAY)

add_text(s, "Thank You!", 0, 7.0, 13.33, 0.45, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────
out = r"c:\Users\hp\Desktop\InsuranceSystem_Presentation.pptx"
prs.save(out)
print("Saved -> " + out)
